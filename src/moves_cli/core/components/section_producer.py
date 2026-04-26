from importlib.resources import files
from pathlib import Path
from typing import Callable, Literal
import contextlib
import io

import instructor

# Suppress pymupdf4llm import warning about layout package
with contextlib.redirect_stdout(io.StringIO()):
    import pymupdf4llm
from docx import Document
from litellm import completion, responses
from pptx import Presentation
from pydantic import BaseModel, Field

from moves_cli.models import Section


class SectionProducer:
    @staticmethod
    def _resolve_temperature(llm_model: str) -> float:
        model = llm_model.strip().lower()
        if "gpt-5" in model:
            return 1.0
        return 0.2

    @staticmethod
    def _build_messages(
        system_prompt: str, presentation_data: str, transcript_data: str
    ) -> list[dict[str, str]]:
        return [
            {"role": "system", "content": system_prompt},
            {
                "role": "user",
                "content": f"Presentation: {presentation_data}\nTranscript: {transcript_data}",
            },
        ]

    @staticmethod
    def _validate_llm_format(llm_format: str) -> Literal["chat", "responses", "auto"]:
        normalized = llm_format.strip().lower()
        if normalized not in {"chat", "responses", "auto"}:
            raise ValueError(
                f"Unsupported LLM format: {llm_format}. "
                "Supported formats: chat, responses, auto"
            )
        return normalized  # type: ignore[return-value]

    @staticmethod
    def _extract_text_from_responses(response: object) -> str:
        output_text = getattr(response, "output_text", None)
        if isinstance(output_text, str) and output_text.strip():
            return output_text

        output_items = None
        if isinstance(response, dict):
            output_items = response.get("output")
        else:
            output_items = getattr(response, "output", None)

        if not isinstance(output_items, list):
            raise ValueError("Responses API returned unexpected output structure")

        text_parts: list[str] = []
        for item in output_items:
            item_type = (
                item.get("type")
                if isinstance(item, dict)
                else getattr(item, "type", None)
            )
            if item_type != "message":
                continue

            content_parts = (
                item.get("content")
                if isinstance(item, dict)
                else getattr(item, "content", None)
            )
            if not isinstance(content_parts, list):
                continue

            for part in content_parts:
                part_type = (
                    part.get("type")
                    if isinstance(part, dict)
                    else getattr(part, "type", None)
                )
                part_text = (
                    part.get("text")
                    if isinstance(part, dict)
                    else getattr(part, "text", None)
                )

                if part_type in {"output_text", "text"} and isinstance(part_text, str):
                    text_parts.append(part_text)

        result = "\n".join(part.strip() for part in text_parts if part and part.strip())
        if result:
            return result

        raise ValueError("Responses API returned empty text output")

    @staticmethod
    def _build_sections_output_model(slide_count: int) -> type[BaseModel]:
        # define output model with schema to reliable extract sections from llm response
        class SectionsOutputModel(BaseModel):
            class SectionItem(BaseModel):
                section_index: int = Field(
                    ...,
                    ge=1,
                    description="Index starting from 1",  # descriptions for llm to understand the schema
                )
                content: str = Field(..., description="Content of the section")

            sections: list[SectionItem] = Field(  # type: ignore
                ...,
                description="List of section items, one for each slide",
                min_items=slide_count,  # must be exact number of slides, min or max.
                max_items=slide_count,
            )

        return SectionsOutputModel

    @staticmethod
    def _build_connection_kwargs(
        llm_base_url: str | None,
        target: Literal["chat", "responses"],
    ) -> dict[str, str]:
        if not llm_base_url:
            return {}

        base_url = llm_base_url.strip()
        if not base_url:
            return {}

        if target == "chat":
            return {"base_url": base_url}

        # litellm.responses() doesn't expose base_url directly, but accepts
        # provider connection params through GenericLiteLLMParams kwargs.
        return {"api_base": base_url}

    def _extract_pdf(
        self, file_path: Path, extraction_type: Literal["transcript", "presentation"]
    ) -> str:
        """Extract text from PDF using PyMuPDF4LLM (optimized for LLM processing).

        Uses markdown conversion which preserves structure better than plain text.
        """
        # PyMuPDF4LLM works with file path (str)
        # Using page_chunks for better structure preservation
        chunks = pymupdf4llm.to_markdown(
            str(file_path),
            page_chunks=True,  # Returns list of dicts, one per page
        )

        match extraction_type:
            case "transcript":
                # Extract all text from all pages, concatenate into one line
                full_text = " ".join(chunk["text"] for chunk in chunks)  # type: ignore
                # Remove extra spaces and newlines for transcript
                result = " ".join(full_text.split())
                return result

            case "presentation":
                # For page-by-page presentation, keep markdown structure
                markdown_sections = []
                for i, chunk in enumerate(chunks):
                    # Clean up markdown text but preserve basic structure
                    page_text = chunk["text"].strip()  # type: ignore
                    cleaned_text = " ".join(page_text.split())
                    markdown_sections.append(f"# Slide Page {i + 1}\n{cleaned_text}")

                return "\n\n".join(markdown_sections)

    def _extract_docx(
        self, file_path: Path, extraction_type: Literal["transcript", "presentation"]
    ) -> str:
        """Extract text from DOCX using python-docx (free, no PyMuPDF Pro needed)."""
        # Read document (python-docx accepts str path)
        doc = Document(str(file_path))

        match extraction_type:
            case "transcript":
                # Extract all text from all paragraphs
                full_text = " ".join(paragraph.text for paragraph in doc.paragraphs)
                result = " ".join(full_text.split())
                return result

            case "presentation":
                # Treat each paragraph as a potential slide
                # This is a heuristic - DOCX doesn't have explicit slides
                markdown_sections = []
                for i, paragraph in enumerate(doc.paragraphs):
                    if paragraph.text.strip():  # Skip empty paragraphs
                        cleaned_text = " ".join(paragraph.text.split())
                        markdown_sections.append(
                            f"# Slide Page {i + 1}\n{cleaned_text}"
                        )

                return (
                    "\n\n".join(markdown_sections)
                    if markdown_sections
                    else "# Slide Page 1\n"
                )

    def _extract_pptx(
        self, file_path: Path, extraction_type: Literal["transcript", "presentation"]
    ) -> str:
        """Extract text from PPTX using python-pptx (free, no PyMuPDF Pro needed)."""
        # Load presentation (python-pptx accepts str path)
        prs = Presentation(str(file_path))

        match extraction_type:
            case "transcript":
                # Extract all text from all slides
                all_text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for paragraph in shape.text_frame.paragraphs:  # type: ignore
                                for run in paragraph.runs:
                                    all_text.append(run.text)

                full_text = " ".join(all_text)
                result = " ".join(full_text.split())
                return result

            case "presentation":
                # Extract text slide by slide
                markdown_sections = []
                for i, slide in enumerate(prs.slides):
                    slide_text_parts = []
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for paragraph in shape.text_frame.paragraphs:  # type: ignore
                                for run in paragraph.runs:
                                    slide_text_parts.append(run.text)

                    slide_text = " ".join(slide_text_parts)
                    cleaned_text = " ".join(slide_text.split())
                    markdown_sections.append(f"# Slide Page {i + 1}\n{cleaned_text}")

                return "\n\n".join(markdown_sections)

    def _extract_txt(
        self, file_path: Path, extraction_type: Literal["transcript", "presentation"]
    ) -> str:
        """Extract text from plain text files."""
        # Read as UTF-8 text
        content = file_path.read_text(encoding="utf-8")

        match extraction_type:
            case "transcript":
                # Single line, no extra spaces
                result = " ".join(content.split())
                return result

            case "presentation":
                # For TXT, treat each line as a potential section
                lines = [line.strip() for line in content.splitlines() if line.strip()]
                markdown_sections = [
                    f"# Slide Page {i + 1}\n{line}" for i, line in enumerate(lines)
                ]
                return (
                    "\n\n".join(markdown_sections)
                    if markdown_sections
                    else "# Slide Page 1\n"
                )

    def _extract_document(
        self, file_path: Path, extraction_type: Literal["transcript", "presentation"]
    ) -> str:
        """Extract text from document.

        Supports: PDF (pymupdf4llm), DOCX (python-docx), PPTX (python-pptx), TXT.
        All formats are free and don't require PyMuPDF Pro.
        """
        suffix = file_path.suffix.lower()

        try:
            match suffix:
                case ".pdf":
                    return self._extract_pdf(file_path, extraction_type)
                case ".docx":
                    return self._extract_docx(file_path, extraction_type)
                case ".pptx":
                    return self._extract_pptx(file_path, extraction_type)
                case ".txt":
                    return self._extract_txt(file_path, extraction_type)
                case _:
                    raise ValueError(
                        f"Unsupported file format: {suffix}. "
                        f"Supported formats: .pdf, .docx, .pptx, .txt"
                    )
        except Exception as e:
            raise RuntimeError(
                f"Document extraction failed for {file_path} ({extraction_type}): {e}"
            ) from e

    def generate_template(self, presentation_path: Path) -> list[Section]:
        """
        Extract slide count from presentation file and generate empty sections.
        Supports: PDF, DOCX, PPTX, TXT (all free, no PyMuPDF Pro needed).
        No LLM call, fully offline.
        """
        suffix = presentation_path.suffix.lower()

        try:
            match suffix:
                case ".pdf":
                    # Use pymupdf4llm to get page chunks
                    chunks = pymupdf4llm.to_markdown(
                        str(presentation_path), page_chunks=True
                    )
                    slide_count = len(chunks)

                case ".docx":
                    # Count non-empty paragraphs as slides
                    doc = Document(str(presentation_path))
                    slide_count = sum(1 for p in doc.paragraphs if p.text.strip())
                    slide_count = max(slide_count, 1)  # At least 1 slide

                case ".pptx":
                    # Count slides directly
                    prs = Presentation(str(presentation_path))
                    slide_count = len(prs.slides)

                case ".txt":
                    # Count non-empty lines as slides
                    content = presentation_path.read_text(encoding="utf-8")
                    lines = [
                        line.strip() for line in content.splitlines() if line.strip()
                    ]
                    slide_count = max(len(lines), 1)  # At least 1 slide

                case _:
                    raise ValueError(
                        f"Unsupported file format: {suffix}. "
                        f"Supported formats: .pdf, .docx, .pptx, .txt"
                    )

            return [
                Section(section_index=i + 1, content="") for i in range(slide_count)
            ]
        except Exception as e:
            raise RuntimeError(
                f"Failed to generate template from {presentation_path}: {e}"
            ) from e

    def estimate_for_files(
        self,
        presentation_path: Path,
        transcript_path: Path,
        llm_model: str,
    ) -> tuple[int, int, float | None]:
        """
        Estimate token count and cost for given files without making LLM call.

        Returns:
            tuple: (slide_count, token_count, estimated_cost_usd or None)
        """
        from litellm import cost_per_token, token_counter

        # Extract data from documents (PDF or TXT)
        presentation_data = self._extract_document(presentation_path, "presentation")
        transcript_data = self._extract_document(transcript_path, "transcript")

        # Count slides
        slide_count = len(presentation_data.split("\n\n"))

        # Build messages for token counting
        system_prompt = (
            files("moves_cli.data").joinpath("llm_instruction.md").read_text()
        )
        messages = self._build_messages(
            system_prompt, presentation_data, transcript_data
        )

        # Count tokens (local, free)
        token_count = token_counter(model=llm_model, messages=messages)

        # Estimate cost (local lookup, free)
        try:
            prompt_cost, _ = cost_per_token(
                model=llm_model, prompt_tokens=token_count, completion_tokens=0
            )
        except Exception:
            prompt_cost = None  # Model pricing not available

        return slide_count, token_count, prompt_cost

    def _call_llm_chat(
        self,
        presentation_data: str,
        transcript_data: str,
        llm_model: str,
        llm_api_key: str,
        llm_base_url: str | None,
    ) -> list[str]:
        slide_count = len(presentation_data.split("\n\n"))
        sections_output_model = self._build_sections_output_model(slide_count)

        try:
            import warnings

            # silence harmless pydantic serialization warnings from litellm/instructor
            # see: https://github.com/BerriAI/litellm/issues/11759
            warnings.filterwarnings(
                "ignore",
                message="Pydantic serializer warnings",
                category=UserWarning,
            )

            # hmm, i need to rewrite this system prompt for broader use cases, current one is too restrictive
            system_prompt = (
                files("moves_cli.data").joinpath("llm_instruction.md").read_text()
            )
            messages = self._build_messages(
                system_prompt, presentation_data, transcript_data
            )
            # we're pathching the litellm with instructor to use any llm with any schema
            client = instructor.from_litellm(
                completion, mode=instructor.Mode.JSON
            )  # json for better llm output quality, also afaik the instructor retries on failure for these
            connection_kwargs = self._build_connection_kwargs(
                llm_base_url=llm_base_url,
                target="chat",
            )
            temperature = self._resolve_temperature(llm_model)

            response = client.chat.completions.create(
                model=llm_model,
                api_key=llm_api_key,
                messages=messages,
                response_model=sections_output_model,
                temperature=temperature,
                reasoning_effort="none",  # we don't need complex reasoning, we just want the sections extracted reliably, so we can set it to none for faster responses
                drop_params=True,  # drop unsupported params instead of raising exceptions
                **connection_kwargs,
            )
            result = [item.content for item in response.sections]
            return result
        except Exception as e:
            raise RuntimeError(f"LLM call failed: {e}") from e

    def _call_llm_responses(
        self,
        presentation_data: str,
        transcript_data: str,
        llm_model: str,
        llm_api_key: str,
        llm_base_url: str | None,
    ) -> list[str]:
        slide_count = len(presentation_data.split("\n\n"))
        sections_output_model = self._build_sections_output_model(slide_count)

        try:
            system_prompt = (
                files("moves_cli.data").joinpath("llm_instruction.md").read_text()
            )
            messages = self._build_messages(
                system_prompt, presentation_data, transcript_data
            )
            connection_kwargs = self._build_connection_kwargs(
                llm_base_url=llm_base_url,
                target="responses",
            )
            temperature = self._resolve_temperature(llm_model)

            response = responses(
                model=llm_model,
                api_key=llm_api_key,
                input=messages,
                text={
                    "format": {
                        "type": "json_schema",
                        "name": "sections_output",
                        "schema": sections_output_model.model_json_schema(),
                        "strict": True,
                    }
                },
                temperature=temperature,
                reasoning={
                    "effort": "none"
                },  # we don't need complex reasoning, we just want the sections extracted reliably, so we can set it to none for faster responses
                drop_params=True,  # drop unsupported params instead of raising exceptions
                **connection_kwargs,
            )

            response_text = self._extract_text_from_responses(response)
            parsed_response = sections_output_model.model_validate_json(response_text)
            return [item.content for item in parsed_response.sections]
        except Exception as e:
            raise RuntimeError(f"LLM call failed: {e}") from e

    def _call_llm(
        self,
        presentation_data: str,
        transcript_data: str,
        llm_model: str,
        llm_api_key: str,
        llm_format: str,
        llm_base_url: str | None,
    ) -> list[str]:
        selected_format = self._validate_llm_format(llm_format)

        if selected_format == "chat":
            return self._call_llm_chat(
                presentation_data=presentation_data,
                transcript_data=transcript_data,
                llm_model=llm_model,
                llm_api_key=llm_api_key,
                llm_base_url=llm_base_url,
            )

        if selected_format == "responses":
            return self._call_llm_responses(
                presentation_data=presentation_data,
                transcript_data=transcript_data,
                llm_model=llm_model,
                llm_api_key=llm_api_key,
                llm_base_url=llm_base_url,
            )

        try:
            return self._call_llm_responses(
                presentation_data=presentation_data,
                transcript_data=transcript_data,
                llm_model=llm_model,
                llm_api_key=llm_api_key,
                llm_base_url=llm_base_url,
            )
        except Exception:
            return self._call_llm_chat(
                presentation_data=presentation_data,
                transcript_data=transcript_data,
                llm_model=llm_model,
                llm_api_key=llm_api_key,
                llm_base_url=llm_base_url,
            )

    def convert_to_markdown(self, sections: list[Section]) -> str:
        """Convert sections to markdown format.

        Format:
            # 1. Slide

            Section content here...

            # 2. Slide

            Another section content...
        """
        lines: list[str] = []
        for section in sections:
            lines.append(f"# {section.section_index}. Slide")
            lines.append("")
            if section.content.strip():
                lines.append(section.content.strip())
                lines.append("")

        return "\n".join(lines)

    def load_from_markdown(self, markdown_content: str) -> list[Section]:
        """Load sections from markdown format.

        Parses section headings as section indices, content follows until next heading.

        Supported heading formats:
        - `# N. Slide` (canonical)
        - `# Slide N` (legacy)
        """
        import re

        sections: list[Section] = []
        heading_pattern = re.compile(
            r"^#\s*(?:(\d+)\.\s*Slide|Slide\s+(\d+))\s*$",
            flags=re.MULTILINE,
        )
        matches = list(heading_pattern.finditer(markdown_content))

        for i, match in enumerate(matches):
            section_index_str = match.group(1) or match.group(2)
            if not section_index_str:
                continue

            section_index = int(section_index_str)
            content_start = match.end()
            content_end = (
                matches[i + 1].start()
                if i + 1 < len(matches)
                else len(markdown_content)
            )
            content = markdown_content[content_start:content_end].strip()
            sections.append(Section(section_index=section_index, content=content))

        return sections

    def generate_sections(
        self,
        presentation_path: Path,
        transcript_path: Path,
        llm_model: str,
        llm_api_key: str,
        llm_format: str,
        llm_base_url: str | None,
        callback: Callable[[str], None] | None = None,
    ) -> list[Section]:
        if callback:
            callback("Extracting presentation data...")
        presentation_data = self._extract_document(presentation_path, "presentation")

        if callback:
            callback("Extracting transcript data...")
        transcript_data = self._extract_document(transcript_path, "transcript")

        if callback:
            callback("Calling LLM...")
        section_contents = self._call_llm(
            presentation_data=presentation_data,
            transcript_data=transcript_data,
            llm_model=llm_model,
            llm_api_key=llm_api_key,
            llm_format=llm_format,
            llm_base_url=llm_base_url,
        )

        generated_sections: list[Section] = []

        for idx, content in enumerate(section_contents):
            section = Section(
                content=content,
                section_index=idx + 1,
            )
            generated_sections.append(section)

        return generated_sections
